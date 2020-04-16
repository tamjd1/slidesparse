from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
from io import StringIO
import os
from nltk.tokenize import sent_tokenize
from pptx import Presentation
import glob


def __find(text, keyword, filename=None, page=0):
    sentences = sent_tokenize(text)
    for s in sentences:
        if keyword.lower() in s.lower():
            print()
            print("- {} --- {}".format(page, filename))
            print()
            for _s in s.split("\n"):
                print("\t {}".format(_s))


def find_in_pptx(path, keyword):  # not tested
    presentation = Presentation(path)
    for i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                __find(shape.text, keyword, os.path.basename(path), i)


def find_in_pdf(path, keyword):
    resource_manager = PDFResourceManager()
    s = StringIO()
    la_params = LAParams()
    device = TextConverter(resource_manager, s, laparams=la_params)
    f = open(path, 'rb')
    interpreter = PDFPageInterpreter(resource_manager, device)

    for page in PDFPage.get_pages(f):
        interpreter.process_page(page)

    text = s.getvalue()
    __find(text, keyword, os.path.basename(path))

    f.close()
    device.close()
    s.close()


def main(keyword):
    slides_dir = "/Users/T/Desktop/slides"  # todo : configure your directory here
    files = glob.glob(os.path.join(slides_dir, "*.pdf"))
    files.extend(glob.glob(os.path.join(slides_dir, "*.ppt")))
    files.extend(glob.glob(os.path.join(slides_dir, "*.pptx")))
    # files = sorted(files, key=lambda x: os.path.getsize(x))
    for f in files:
        print("--- Searching in {} ...".format(f))
        if f.endswith(".pdf"):
            find_in_pdf(f, keyword)
        elif f.endswith(".ppt") or f.endswith(".pptx"):
            find_in_pptx(f, keyword)


if __name__ == '__main__':
    import sys
    args = sys.argv
    __keyword = args[1]
    print("Searching for {} ...".format(__keyword))
    main(__keyword)
